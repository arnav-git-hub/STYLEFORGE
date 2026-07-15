import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation(github_username="arnav-git-hub"):
    prs = Presentation()
    
    # Set slide dimensions to 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Color palette
    COLOR_BG = RGBColor(18, 18, 18)         # Dark charcoal
    COLOR_TEXT = RGBColor(245, 245, 247)    # Near white
    COLOR_PRIMARY = RGBColor(168, 85, 247)  # Cyber purple / Lavender
    COLOR_MUTED = RGBColor(156, 163, 175)   # Muted grey
    COLOR_ACCENT = RGBColor(234, 179, 8)    # Amber/Gold

    blank_layout = prs.slide_layouts[6] # Blank slide layout

    # Helper function to add background color
    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper function to add header/footer
    def add_header_footer(slide, title_text):
        # Header
        txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text.upper()
        p.font.name = 'Georgia'
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        
        # Muted divider line
        divider = slide.shapes.add_shape(1, Inches(0.75), Inches(1.3), Inches(11.833), Inches(0.02))
        divider.fill.solid()
        divider.fill.fore_color.rgb = COLOR_MUTED
        divider.line.color.rgb = COLOR_MUTED

        # Footer
        footerBox = slide.shapes.add_textbox(Inches(0.75), Inches(6.8), Inches(11.833), Inches(0.4))
        ftf = footerBox.text_frame
        fp = ftf.paragraphs[0]
        fp.text = "Edge Hack 2026 Submission  |  Project: DesignAlchemy"
        fp.font.name = 'Arial'
        fp.font.size = Pt(10)
        fp.font.color.rgb = COLOR_MUTED

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide (Dark Theme)
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, COLOR_BG)
    
    # Title & Subtitle in a single text box
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "DESIGNALCHEMY"
    p1.font.name = 'Georgia'
    p1.font.size = Pt(64)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEXT
    p1.alignment = PP_ALIGN.LEFT
    
    p2 = tf1.add_paragraph()
    p2.text = "Bespoke AI-Powered Digital Fashion Studio & Vector Layering Suite"
    p2.font.name = 'Arial'
    p2.font.size = Pt(22)
    p2.font.color.rgb = COLOR_PRIMARY
    p2.space_before = Pt(14)
    p2.alignment = PP_ALIGN.LEFT

    p3 = tf1.add_paragraph()
    p3.text = f"Edge Hack 2026 Submission  |  Built by Team StyleForge"
    p3.font.name = 'Arial'
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLOR_MUTED
    p3.space_before = Pt(40)
    p3.alignment = PP_ALIGN.LEFT

    # -------------------------------------------------------------
    # SLIDE 2: Problem Statement
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, COLOR_BG)
    add_header_footer(slide2, "The Creative Problem")

    content_box2 = slide2.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(4.5))
    tf2 = content_box2.text_frame
    tf2.word_wrap = True

    bullets2 = [
        ("Traditional Prototyping is Slow & Wasteful", 
         "Creating physical fashion swatches and assembling prototype silhouette layers involves high fabric waste and slow iterative turnaround times."),
        ("Lack of Direct Human-AI Collaborative Design Canvas", 
         "Generic image generators do not allow structured, incremental editing of designs, leaving designers unable to adjust fabric textures, layers, or opacities dynamically."),
        ("Decoupled Design Analysis", 
         "Standard tools lack an integrated design co-pilot that programmatically understands the composition of active work on a canvas (drapes, layers, textures) to give real-time styling advice.")
    ]

    for title, desc in bullets2:
        p_title = tf2.add_paragraph()
        p_title.text = "•  " + title
        p_title.font.name = 'Georgia'
        p_title.font.size = Pt(20)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_ACCENT
        p_title.space_before = Pt(18)

        p_desc = tf2.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = 'Arial'
        p_desc.font.size = Pt(15)
        p_desc.font.color.rgb = COLOR_TEXT
        p_desc.space_before = Pt(4)
        p_desc.level = 0

    # -------------------------------------------------------------
    # SLIDE 3: The Solution
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, COLOR_BG)
    add_header_footer(slide3, "Our Solution: DesignAlchemy")

    content_box3 = slide3.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(4.5))
    tf3 = content_box3.text_frame
    tf3.word_wrap = True

    bullets3 = [
        ("Generative Textile Wall", 
         "Weave bespoke fabric textures on-demand from visual/text prompts, instantly generating tileable fabric swatches using Gemini."),
        ("Mannequin Layering Canvas", 
         "A modular digital workspace supporting multi-layer compositions (Base Silhouette, Body Panel, Collar Accents, Lace Embellishments) with precise opacity controls."),
        ("Interactive Studio Muse Co-pilot", 
         "An elite AI design strategist directly connected to the canvas state. It parses the current active swatches, opacity levels, and structural layers to provide contextual couture recommendations.")
    ]

    for title, desc in bullets3:
        p_title = tf3.add_paragraph()
        p_title.text = "•  " + title
        p_title.font.name = 'Georgia'
        p_title.font.size = Pt(20)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_PRIMARY
        p_title.space_before = Pt(18)

        p_desc = tf3.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = 'Arial'
        p_desc.font.size = Pt(15)
        p_desc.font.color.rgb = COLOR_TEXT
        p_desc.space_before = Pt(4)

    # -------------------------------------------------------------
    # SLIDE 4: Key Features & Demonstration Flow
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, COLOR_BG)
    add_header_footer(slide4, "Feature Spotlight & User Flow")

    content_box4 = slide4.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(4.5))
    tf4 = content_box4.text_frame
    tf4.word_wrap = True

    bullets4 = [
        ("Weave Swatches:", "Type prompts like 'Cyberpunk Gold Foil' to instantly generate detailed fashion textures."),
        ("Layer and Mix:", "Select, stack, toggle visibility, and fade fabric swatches onto specific mannequin zones."),
        ("Co-create with Studio Muse:", "Converse with your AI muse. It reads active colors and fabrics, suggesting styling lines, silhouette extension, or runway accessories."),
        ("Offline Mode Simulation:", "If connection fails, the app falls back to a mock mode preserving continuous designer flow.")
    ]

    for label, text in bullets4:
        p = tf4.add_paragraph()
        p.text = "•  "
        p.font.name = 'Arial'
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(14)
        
        run_label = p.add_run()
        run_label.text = label + " "
        run_label.font.bold = True
        run_label.font.color.rgb = COLOR_ACCENT
        
        run_text = p.add_run()
        run_text.text = text
        run_text.font.bold = False
        run_text.font.color.rgb = COLOR_TEXT

    # -------------------------------------------------------------
    # SLIDE 5: Architecture & Technologies
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, COLOR_BG)
    add_header_footer(slide5, "Technical Architecture")

    content_box5 = slide5.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(4.5))
    tf5 = content_box5.text_frame
    tf5.word_wrap = True

    bullets5 = [
        ("Frontend Engine", "Built on Vite + React 19 and TypeScript for an extremely fast, secure, and typed user interface."),
        ("Fluid Workspace Animations", "Utilizes motion/react (Framer Motion) for micro-animations, layer sorting, and tab transitions."),
        ("Backend Services & API Proxy", "Node server proxy with Express handling server-side API requests to protect sensitive keys."),
        ("Gemini Cognitive Layer", "Integrates @google/genai SDK leveraging the gemini-2.0-flash model for high-speed content generation and interactive design chat."),
        ("Vercel Edge Ready", "Fully configured for serverless deployment on Vercel via vercel.json routing and API serverless entry points.")
    ]

    for title, desc in bullets5:
        p_title = tf5.add_paragraph()
        p_title.text = "•  " + title
        p_title.font.name = 'Georgia'
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_PRIMARY
        p_title.space_before = Pt(12)

        p_desc = tf5.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = 'Arial'
        p_desc.font.size = Pt(14)
        p_desc.font.color.rgb = COLOR_TEXT
        p_desc.space_before = Pt(2)

    # -------------------------------------------------------------
    # SLIDE 6: Deliverables & Submission Links
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, COLOR_BG)
    add_header_footer(slide6, "Submission Deliverables")

    content_box6 = slide6.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(4.5))
    tf6 = content_box6.text_frame
    tf6.word_wrap = True

    p_intro = tf6.add_paragraph()
    p_intro.text = "Team StyleForge is pleased to submit the following deliverables for review:"
    p_intro.font.name = 'Arial'
    p_intro.font.size = Pt(16)
    p_intro.font.color.rgb = COLOR_TEXT
    p_intro.space_before = Pt(10)

    links = [
        ("GitHub Repository Link", f"https://github.com/{github_username}/styleforge", "Source code, server structure, layering canvas, and custom serverless functions."),
        ("Interactive Demo Video", f"https://github.com/{github_username}/styleforge/raw/main/artifacts/designalchemy_walkthrough.webp", "Recorded walkthrough showing the complete user flow from fabric generation to Studio Muse interaction."),
        ("Local Dev Environment", "http://localhost:3000", "Live design studio workspace running on the local development port.")
    ]

    for title, link, desc in links:
        p_title = tf6.add_paragraph()
        p_title.text = "•  " + title
        p_title.font.name = 'Georgia'
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_ACCENT
        p_title.space_before = Pt(18)

        p_link = tf6.add_paragraph()
        p_link.text = f"URL: {link}"
        p_link.font.name = 'Arial'
        p_link.font.size = Pt(15)
        p_link.font.bold = True
        p_link.font.color.rgb = COLOR_PRIMARY
        p_link.space_before = Pt(2)

        p_desc = tf6.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = 'Arial'
        p_desc.font.size = Pt(14)
        p_desc.font.color.rgb = COLOR_MUTED
        p_desc.space_before = Pt(2)

    prs.save("e:\\styleforge\\edge_hack_presentation.pptx")
    print("Presentation saved successfully at e:\\styleforge\\edge_hack_presentation.pptx")

if __name__ == "__main__":
    github_user = sys.argv[1] if len(sys.argv) > 1 else "arnav-git-hub"
    create_presentation(github_user)
